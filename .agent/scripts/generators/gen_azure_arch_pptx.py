"""
gen_azure_arch_pptx.py  v2  --  Diagrama de fluxo estilo GitOps / Azure Architecture
Modelo: background branco, icones Azure style, conectores com setas,
        badges numerados verdes, whitespace generoso (igual imagem de referencia).
5 slides: Capa | Fluxo Principal | Stack Azure | Fluxo DevOps CI-CD | Roadmap
Saida: _outputs/kare-azure/upstream/KARE_Azure_Architecture.pptx
"""

import os
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Paleta ────────────────────────────────────────────────────────────────────
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_F5       = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_LIGHT    = RGBColor(0xE0, 0xE0, 0xE0)
GRAY_MED      = RGBColor(0xA0, 0xA0, 0xA0)
GRAY_DARK     = RGBColor(0x50, 0x50, 0x50)
TEXT_DARK     = RGBColor(0x1A, 0x1A, 0x2E)
AZ_NAVY       = RGBColor(0x00, 0x20, 0x50)
AZ_BLUE       = RGBColor(0x00, 0x78, 0xD4)
AZ_BLUE_LT    = RGBColor(0xD6, 0xEB, 0xFF)
AZ_PURPLE     = RGBColor(0x77, 0x3A, 0xDC)
AZ_PURPLE_LT  = RGBColor(0xED, 0xE8, 0xFB)
AZ_TEAL       = RGBColor(0x00, 0x7A, 0x78)
AZ_TEAL_LT    = RGBColor(0xCE, 0xEE, 0xEE)
AZ_GREEN      = RGBColor(0x10, 0x79, 0x3B)
AZ_GREEN_LT   = RGBColor(0xD8, 0xF3, 0xE5)
AZ_RED        = RGBColor(0xC0, 0x37, 0x26)
AZ_RED_LT     = RGBColor(0xFB, 0xE3, 0xE1)
AZ_ORANGE     = RGBColor(0xEF, 0x6C, 0x00)
AZ_ORG_LT     = RGBColor(0xFD, 0xEF, 0xE0)
GOLD          = RGBColor(0xFF, 0xB9, 0x00)
GOLD_DK       = RGBColor(0xB0, 0x70, 0x00)
ACCENT_PUR    = RGBColor(0x66, 0x00, 0x99)
ACCENT_MAG    = RGBColor(0xBB, 0x00, 0x7D)
CYAN_ACC      = RGBColor(0x50, 0xE6, 0xFF)
P_BLUE        = RGBColor(0x00, 0x78, 0xD4)
P_BLUE2       = RGBColor(0x50, 0xA8, 0xE0)
P_BLUE3       = RGBColor(0x90, 0xC8, 0xF0)


# ── Primitivos ────────────────────────────────────────────────────────────────

def _rect(sl, x, y, w, h, fill, border=None, bpt=1.0, rx=False):
    stype = 5 if rx else 1
    s = sl.shapes.add_shape(stype, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(bpt)
    else:
        s.line.fill.background()
    return s


def _oval(sl, x, y, w, h, fill, border=None, bpt=1.0):
    s = sl.shapes.add_shape(9, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(bpt)
    else:
        s.line.fill.background()
    return s


def _label(sl, text, x, y, w, h, size=9, bold=False, color=TEXT_DARK,
           align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Segoe UI"
    return tb


def _label_in(shape, line1, line2=None, size1=8, size2=6.5, bold=True, col1=AZ_NAVY):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = line1
    r.font.size = Pt(size1)
    r.font.bold = bold
    r.font.color.rgb = col1
    r.font.name = "Segoe UI"
    if line2:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = line2
        r2.font.size = Pt(size2)
        r2.font.color.rgb = GRAY_DARK
        r2.font.name = "Segoe UI"


def _arrow(sl, x1, y1, x2, y2, color=GRAY_DARK, pt=1.4):
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
    conn = sl.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(pt)
    spPr = conn._element.spPr
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        te = etree.SubElement(ln, qn('a:tailEnd'))
        te.set('type', 'arrow')
        te.set('w', 'med')
        te.set('len', 'med')
    return conn


def _badge(sl, x, y, num, color=AZ_GREEN, sz=0.3):
    c = _oval(sl, x, y, sz, sz, color)
    _label_in(c, str(num), size1=9, bold=True, col1=WHITE)


def _people(sl, x, y, n=3, base_color=P_BLUE):
    """Overlapping silhouettes like the reference image"""
    shades = [P_BLUE3, P_BLUE2, P_BLUE]
    for i in range(min(n, 3)):
        ox = x + i * 0.2
        oy = y - i * 0.05
        c = shades[i % 3] if base_color == P_BLUE else base_color
        _oval(sl, ox + 0.12, oy, 0.24, 0.24, c)
        _rect(sl, ox + 0.04, oy + 0.26, 0.40, 0.22, c, rx=True)


def _person(sl, x, y, color=AZ_BLUE):
    _oval(sl, x + 0.18, y, 0.28, 0.28, color)
    _rect(sl, x + 0.06, y + 0.30, 0.52, 0.26, color, rx=True)


# ── Header/Footer ─────────────────────────────────────────────────────────────

def _hdr(sl, title, subtitle=None, light=False):
    bg = GRAY_F5 if light else AZ_NAVY
    tc = AZ_NAVY if light else WHITE
    _rect(sl, 0, 0, 13.333, 0.55, bg)
    _rect(sl, 0, 0, 0.07, 0.55, AZ_BLUE)
    _label(sl, title, 0.22, 0.04, 11.5, 0.34, size=16, bold=True, color=tc)
    if subtitle:
        _label(sl, subtitle, 0.22, 0.37, 11.5, 0.17,
               size=7, color=GRAY_DARK if light else GRAY_LIGHT)
    _rect(sl, 12.08, 0.07, 1.15, 0.40, ACCENT_PUR, rx=True)
    _label(sl, "KARE", 12.08, 0.07, 1.15, 0.40,
           size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def _ftr(sl, txt=None):
    _rect(sl, 0, 7.22, 13.333, 0.28, AZ_NAVY)
    msg = txt or "KARE Project  |  PI Planning CLOCK02 26  |  KARE Agile Agent  |  Maio 2026"
    _label(sl, msg, 0.3, 7.24, 12.5, 0.22, size=6.5, color=GRAY_LIGHT)


# ── Node helpers ──────────────────────────────────────────────────────────────

def _node(sl, x, y, w, h, line1, line2=None, fill=AZ_BLUE_LT, border=AZ_BLUE,
          bpt=2.0, name_size=8.5, sub_size=6.5):
    """Rounded-rect icon node (Azure style)"""
    s = _rect(sl, x, y, w, h, fill, border, bpt, rx=True)
    _label_in(s, line1, line2, size1=name_size, size2=sub_size, bold=True, col1=border)
    return s


def _kv_icon(sl, x, y, label_txt, step_n, badge_col=AZ_GREEN):
    """Key Vault icon: gold circle + label + step badge"""
    shp = _oval(sl, x, y, 1.1, 1.1, GOLD, GOLD_DK, 2.5)
    _label_in(shp, "KEY", "VAULT", size1=10, size2=8, bold=True, col1=GOLD_DK)
    _label(sl, label_txt, x - 0.2, y + 1.14, 1.5, 0.22,
           size=8, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    _badge(sl, x + 0.88, y - 0.07, step_n, color=badge_col)


def _m365_icon(sl, x, y, w, h, label_top, label_bot, fill=AZ_BLUE_LT, border=AZ_BLUE, step_n=None):
    """Microsoft 365 / service icon (like orange square in reference)"""
    # Outer colored box
    s = _rect(sl, x, y, w, h, fill, border, 2.5, rx=False)
    _label_in(s, label_top, label_bot, size1=8, size2=6.5, bold=True, col1=border)
    if step_n:
        _badge(sl, x + w - 0.18, y - 0.16, step_n)
    return s


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — CAPA
# ═══════════════════════════════════════════════════════════════════════════════

def slide_capa(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, 13.333, 7.5, AZ_NAVY)
    _rect(sl, 0, 0, 0.5, 7.5, AZ_BLUE)
    _rect(sl, 0.5, 6.88, 12.833, 0.1, ACCENT_MAG)
    # decorative arc top-right
    _oval(sl, 9.8, -1.5, 5.0, 5.0, RGBColor(0x00, 0x38, 0x80))

    _label(sl, "KARE Agile Agent", 0.85, 1.4, 11.0, 1.4,
           size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    _label(sl, "Arquitetura Azure + Microsoft Teams",
           0.85, 3.0, 11.0, 0.75, size=26, color=CYAN_ACC, align=PP_ALIGN.LEFT)
    _label(sl, "Planejamento de infraestrutura para disponibilizacao corporativa em nuvem",
           0.85, 3.95, 11.0, 0.5, size=13, color=RGBColor(0xCC, 0xCC, 0xFF))

    tags = [
        ("34 Agentes IA",      AZ_PURPLE),
        ("65+ Skills",         AZ_TEAL),
        ("47 Slash Commands",  AZ_GREEN),
        ("3 MCP Servers",      AZ_ORANGE),
    ]
    for i, (t, c) in enumerate(tags):
        tx = 0.85 + i * 2.52
        _rect(sl, tx, 4.85, 2.28, 0.5, c, rx=True)
        _label(sl, t, tx, 4.85, 2.28, 0.5,
               size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _label(sl, "KARE Project  |  PI Planning CLOCK02 26  |  Cliente  |  Maio 2026",
           0.85, 7.05, 11.5, 0.3, size=9, color=GRAY_MED)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — FLUXO PRINCIPAL (estilo referencia GitOps)
# ═══════════════════════════════════════════════════════════════════════════════
#
#  Estrutura espelhando a imagem de referencia:
#
#   [Azure DevOps logo]          [Key Vault]           [Validate] → [M365 Stage]
#                                     ↑ 6
#  [Fork] →3→ [Main Repo] →5→ [Pipeline] → [M365DSC] →7→     ↓ Deploy
#    ↑ 2                   4↓                              [Approve] → [M365 Prod]
#  [Config]←1← [Admin]     [Pipeline]     [Admins]
#                           check PR        review
#
#  Para KARE:
#   [KARE logo]            [Key Vault]              [Squads] → [Conf Stage]
#                               ↑ 6
#  [Channel] →3→ [Orch] →5→ [AI Foundry] → [MCP] →7→     ↓ Deploy
#    ↑ 2              4↓                          [Admins] → [Conf Prod]
#  [Config]←1← [Squad]   [Bot/APIM check]  [Squads] review

def slide_fluxo(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    # white background
    bk = sl.background
    bk.fill.solid()
    bk.fill.fore_color.rgb = WHITE

    _hdr(sl, "Fluxo de Operacao KARE  --  GitOps-IA Pattern no Azure",
         "Como um /slash-command percorre toda a infraestrutura Azure do KARE Agile Agent",
         light=True)
    _ftr(sl)

    # ── coluna de fundo colorido (muito suave, como o reference) ──────────────
    col_bands = [
        (0.08,  2.12,  GRAY_F5),
        (2.20,  2.08,  RGBColor(0xFD, 0xF4, 0xF3)),
        (4.28,  5.25,  RGBColor(0xF0, 0xEB, 0xFD)),
        (9.53,  3.72,  RGBColor(0xD0, 0xEF, 0xEF)),
    ]
    for cx, cw, cbg in col_bands:
        _rect(sl, cx, 0.55, cw, 6.40, cbg)

    # ── KARE Logo decorativo (top-left, como Azure DevOps logo) ──────────────
    logo = _rect(sl, 0.15, 0.65, 1.85, 1.4, AZ_BLUE_LT, AZ_BLUE, 2.5, rx=True)
    _label_in(logo, "KARE", "Agile Agent", size1=18, size2=10, bold=True, col1=AZ_BLUE)

    # ── KEY VAULT (top center) ────────────────────────────────────────────────
    kv_x, kv_y = 5.75, 0.66
    _kv_icon(sl, kv_x, kv_y,
             "Key Vaults", step_n=6)
    _label(sl, "Get credentials used to\ninvoke AI agents", kv_x - 0.35, kv_y + 1.38, 1.8, 0.38,
           size=6.5, italic=True, color=GRAY_DARK, align=PP_ALIGN.CENTER, wrap=True)

    # ── LEFT SECTION: Squad / Config / Channel ────────────────────────────────
    # Squad person (bottom-left)
    _person(sl, 0.35, 4.85)
    _label(sl, "Admin / Squad", 0.1, 5.55, 1.6, 0.22,
           size=8, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    _badge(sl, 1.52, 4.82, 1)
    _label(sl, "Add / Modify /\nDelete config", 0.1, 5.8, 1.6, 0.32,
           size=7, italic=True, color=GRAY_DARK, align=PP_ALIGN.CENTER)

    # Config file (like document icon in reference)
    _rect(sl, 0.25, 3.52, 1.4, 1.0, WHITE, GRAY_MED, 1.5)
    _rect(sl, 1.38, 3.52, 0.27, 0.27, GRAY_LIGHT)  # fold corner
    _label(sl, "BACKLOG.md\nSTORY.md\nConfig files", 0.3, 3.58, 1.3, 0.82,
           size=7, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    _label(sl, "Config file", 0.2, 4.57, 1.6, 0.22,
           size=8, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # Teams Channel Fork (like "Fork Admin 1" diamond in reference)
    fork = _node(sl, 0.2, 2.18, 1.85, 0.82,
                 "Teams Channel", "/slash-command",
                 AZ_BLUE_LT, AZ_BLUE, 2.5)
    _badge(sl, 1.75, 2.15, 2)
    _label(sl, "Commit &\nSync", 1.0, 3.07, 0.9, 0.38,
           size=7, italic=True, color=GRAY_DARK, align=PP_ALIGN.CENTER)

    # arrows left column
    _arrow(sl, 0.95, 4.85, 0.95, 4.52, GRAY_MED, pt=1.2)   # person -> config (up)
    _arrow(sl, 0.95, 3.52, 0.95, 3.0, GRAY_MED, pt=1.2)    # config -> channel (up)

    # ── MIDDLE-LEFT: Bot Service + APIM (como "Main Repository + Pipeline") ───
    # Bot Service (like "Main repository" diamond)
    main_repo = _node(sl, 2.28, 2.18, 1.8, 0.82,
                      "Azure AI Bot", "Teams Direct Line",
                      AZ_BLUE_LT, AZ_BLUE, 2.5)
    _badge(sl, 3.78, 2.15, 3)
    _label(sl, "Create\nSlash-Command", 2.48, 1.8, 1.4, 0.34,
           size=7, italic=True, color=GRAY_DARK, align=PP_ALIGN.CENTER)

    # Pipeline check PR (like the "Pipeline" icon below main repo)
    pipe_check = _node(sl, 2.28, 3.52, 1.8, 0.72,
                       "Bot Service", "Check command",
                       AZ_RED_LT, AZ_RED, 1.5)
    _badge(sl, 3.78, 3.5, 4)

    # Admins (code review group) -- like "Admins" people group in reference
    _people(sl, 2.5, 4.65, n=3)
    _label(sl, "Admins / Squad", 2.35, 5.42, 1.8, 0.22,
           size=8, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    _badge(sl, 3.85, 4.62, 5)
    _label(sl, "Code Review &\nMerge Command", 2.38, 5.66, 1.72, 0.32,
           size=7, italic=True, color=GRAY_DARK, align=PP_ALIGN.CENTER)

    _arrow(sl, 3.18, 3.0, 3.18, 3.52, GRAY_MED, pt=1.2)   # bot -> pipeline
    _arrow(sl, 3.18, 4.24, 3.18, 4.65, GRAY_MED, pt=1.2)  # pipeline -> admins

    # ── CENTER: KARE Orchestrator + AI Foundry (como Multi-stage Pipeline + M365DSC) ──
    # Multi-stage Pipeline equivalent = KARE Orchestrator
    orch = _node(sl, 4.38, 2.18, 2.0, 0.82,
                 "KARE Orchestrator", "AI Foundry Master",
                 AZ_PURPLE_LT, AZ_PURPLE, 2.5)
    _label(sl, "Multi-stage\nPipeline", 4.55, 1.8, 1.65, 0.34,
           size=7, italic=True, color=GRAY_DARK, align=PP_ALIGN.CENTER)

    # Microsoft365DSC equivalent = AI Foundry Agents
    dsc = _node(sl, 6.68, 2.18, 2.0, 0.82,
                "AI Foundry", "34 Agentes exec.",
                AZ_PURPLE_LT, AZ_PURPLE, 2.5)
    _badge(sl, 8.37, 2.15, 7)
    _label(sl, "Deploys changes from\nAI tasks via agents", 6.62, 3.05, 2.12, 0.38,
           size=6.5, italic=True, color=GRAY_DARK, align=PP_ALIGN.CENTER, wrap=True)

    # Key Vault arrow down to Orchestrator
    _arrow(sl, kv_x + 0.55, kv_y + 1.1, 5.38, 2.18, GOLD_DK, pt=1.5)

    # Context Engine RAG (below, supporting data layer)
    ctx = _node(sl, 4.38, 3.55, 4.3, 0.72,
                "Context Engine RAG  +  Azure AI Search  +  Cosmos DB",
                "FastAPI container + AI Search hibrido + pgvector",
                AZ_TEAL_LT, AZ_TEAL, 1.5)

    # arrows orchestrator area
    _arrow(sl, 6.38, 2.59, 6.68, 2.59, AZ_PURPLE, pt=2.0)   # orch -> agents
    _arrow(sl, 5.38, 3.0, 5.38, 3.55, AZ_TEAL, pt=1.2)      # orch -> context
    _arrow(sl, 7.68, 3.0, 7.68, 3.55, AZ_TEAL, pt=1.2)      # agents -> context

    # MCP Servers
    mcp = _node(sl, 4.38, 4.55, 4.3, 0.72,
                "MCP Servers  (Atlassian  /  Figma  /  ServiceNow)",
                "3 containers ACA   |   Jira + Confluence + Figma + SNOW",
                AZ_TEAL_LT, AZ_TEAL, 1.5)
    _arrow(sl, 6.53, 4.27, 6.53, 4.55, AZ_TEAL, pt=1.2)

    # ── RIGHT: Deploy targets (M365 Staging / Production) ─────────────────────
    # Admins validate (top-right people)
    _people(sl, 9.72, 1.05, n=3)
    _label(sl, "Squads", 9.6, 1.82, 1.5, 0.22,
           size=8, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    _badge(sl, 11.0, 1.02, 8)
    _label(sl, "Validate\nchanges", 9.6, 2.06, 1.5, 0.28,
           size=7, italic=True, color=GRAY_DARK, align=PP_ALIGN.CENTER)

    # M365 Staging = Confluence Homologacao
    conf_stg = _m365_icon(sl, 11.22, 1.18, 1.95, 1.0,
                          "Confluence", "M365 STAGING",
                          AZ_BLUE_LT, AZ_BLUE)
    _label(sl, "Microsoft 365\nstaging", 11.15, 2.22, 2.1, 0.3,
           size=7.5, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # Admins approve (bottom-right)
    _people(sl, 9.72, 4.1, n=3)
    _label(sl, "Admins", 9.6, 4.87, 1.5, 0.22,
           size=8, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    _badge(sl, 11.0, 4.08, 9)
    _label(sl, "Approve change\nfor deployment", 9.6, 5.12, 1.5, 0.32,
           size=7, italic=True, color=GRAY_DARK, align=PP_ALIGN.CENTER)

    # M365 Production = Confluence + Jira Producao
    conf_prd = _m365_icon(sl, 11.22, 4.28, 1.95, 1.0,
                          "Jira + Conf.", "M365 PRODUCTION",
                          AZ_GREEN_LT, AZ_GREEN)
    _label(sl, "Microsoft 365\nproduction", 11.15, 5.32, 2.1, 0.3,
           size=7.5, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # Arrows validate/approve
    _arrow(sl, 10.55, 1.82, 11.22, 1.68, AZ_BLUE, pt=1.3)
    _label(sl, "Deploy", 11.38, 1.38, 0.7, 0.22,
           size=7, italic=True, color=GRAY_DARK)
    _arrow(sl, 10.55, 4.87, 11.22, 4.78, AZ_GREEN, pt=1.3)
    _label(sl, "Deploy", 11.38, 5.0, 0.7, 0.22,
           size=7, italic=True, color=GRAY_DARK)

    # ── MAIN HORIZONTAL FLOW ARROWS ───────────────────────────────────────────
    # Teams Channel -> Bot Service
    _arrow(sl, 2.05, 2.59, 2.28, 2.59, AZ_BLUE, pt=1.8)
    # Bot Service -> Orchestrator
    _arrow(sl, 4.08, 2.59, 4.38, 2.59, AZ_BLUE, pt=1.8)
    # Agents -> MCP
    _arrow(sl, 8.68, 2.59, 9.18, 2.59, AZ_PURPLE, pt=1.5)
    # MCP -> Stage (diagonal up)
    _arrow(sl, 8.68, 5.0, 11.22, 1.68, AZ_TEAL, pt=1.3)
    # MCP -> Prod (horizontal + slight down)
    _arrow(sl, 8.68, 5.1, 11.22, 4.88, AZ_TEAL, pt=1.3)

    # Labels on arrows
    _label(sl, "Check\nPull Request", 2.1, 1.72, 1.3, 0.32,
           size=6.5, italic=True, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    _label(sl, "Code Review &\nMerge PR", 3.58, 1.72, 1.3, 0.32,
           size=6.5, italic=True, color=GRAY_DARK, align=PP_ALIGN.CENTER)

    # Microsoft Azure badge (bottom left, exact like reference)
    _rect(sl, 0.1, 6.88, 1.55, 0.28, WHITE, GRAY_LIGHT, 0.5)
    _label(sl, "Microsoft Azure", 0.18, 6.9, 1.38, 0.22,
           size=7, bold=True, color=AZ_BLUE)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — STACK AZURE (6 categorias em grade)
# ═══════════════════════════════════════════════════════════════════════════════

def slide_stack(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bk = sl.background
    bk.fill.solid()
    bk.fill.fore_color.rgb = GRAY_F5

    _hdr(sl, "Stack Tecnologica Azure  --  14 Servicos por Categoria Funcional")
    _ftr(sl)

    cats = [
        ("Interface / Bot",       AZ_BLUE,   AZ_BLUE_LT, [
            ("Azure AI Bot Service",        "Teams Direct Line  --  SKU S1 Standard"),
            ("Microsoft Teams App",          "Bot + Tab + 47 slash-commands  --  Adaptive Cards"),
        ]),
        ("Gateway / Seguranca",   AZ_RED,    AZ_RED_LT, [
            ("Azure API Management",         "Rate limiting, autenticacao, versioning  --  Standard v2"),
            ("Azure Front Door Premium",     "WAF, DDoS protection, SSL termination, CDN global"),
            ("Microsoft Entra ID",           "SSO corporativo  --  RBAC por squad e funcao"),
            ("Azure Key Vault",              "Secrets: OpenAI, Atlassian, Figma, ServiceNow"),
        ]),
        ("Inteligencia Artificial", AZ_PURPLE, AZ_PURPLE_LT, [
            ("Azure AI Foundry Agent Svc",   "34 agentes orquestrados  --  paralelo e sequencial"),
            ("Azure OpenAI GPT-4o / o1",     "LLM principal + text-embedding-3-large  --  PTU Brazil"),
            ("Azure AI Search",              "Indice hibrido BM25 + vetorial  --  RAG pipeline"),
        ]),
        ("Compute / Containers",  AZ_ORANGE, AZ_ORG_LT, [
            ("Azure Container Apps",         "Context Engine FastAPI + 3 MCP Servers  --  scale-to-zero"),
            ("Azure Container Registry",     "Registro Docker privado  --  pipeline CI/CD integrado"),
            ("Azure Functions",              "Scripts Python on-demand: PDF, PPTX, ingestao RAG"),
        ]),
        ("Dados / Persistencia",  AZ_GREEN,  AZ_GREEN_LT, [
            ("Azure Cosmos DB",              "Grafo de conhecimento, decisoes, historico  --  NoSQL"),
            ("Azure Blob Storage",           "Artefatos _outputs, uploads, exports  --  LRS / GRS"),
            ("Azure Redis Cache",            "Session cache agentes e RAG  --  TTL 72h"),
            ("Azure PostgreSQL Flex",        "Metadados de projetos e usuarios  --  extensao pgvector"),
        ]),
        ("Observab. / DevOps",    AZ_TEAL,   AZ_TEAL_LT, [
            ("Azure Monitor + App Insights", "Metricas, traces, dashboards executivos, alertas SLO"),
            ("Microsoft Defender for Cloud", "Postura de seguranca, deteccao de ameacas, compliance"),
            ("Azure DevOps Pipelines",       "CI/CD  --  build, test e deploy de containers/Functions"),
        ]),
    ]

    CW = 6.38
    GAP = 0.14
    for idx, (cat, col, bg_c, svcs) in enumerate(cats):
        cx = 0.18 + (idx % 2) * (CW + GAP)
        cy = 0.68 + (idx // 2) * 2.14

        _rect(sl, cx, cy, CW, 0.34, col)
        _label(sl, cat, cx + 0.12, cy + 0.04, CW - 0.2, 0.26,
               size=10, bold=True, color=WHITE)

        body_h = 1.75
        _rect(sl, cx, cy + 0.34, CW, body_h, bg_c, col, 0.5)

        rh = body_h / len(svcs)
        for si, (nm, desc) in enumerate(svcs):
            iy = cy + 0.34 + si * rh
            _oval(sl, cx + 0.1, iy + rh / 2 - 0.07, 0.12, 0.12, col)
            _label(sl, nm, cx + 0.28, iy + 0.05, CW - 0.38, 0.26,
                   size=8.5, bold=True, color=col)
            _label(sl, desc, cx + 0.28, iy + 0.31, CW - 0.38, rh - 0.35,
                   size=6.5, color=GRAY_DARK, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — FLUXO DEVOPS / CI-CD
# ═══════════════════════════════════════════════════════════════════════════════

def slide_devops(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bk = sl.background
    bk.fill.solid()
    bk.fill.fore_color.rgb = WHITE

    _hdr(sl, "Fluxo DevOps / CI-CD  --  Deploy e Operacao do KARE no Azure",
         "Pipeline automatizado: codigo  ->  container  ->  AI Foundry  ->  Teams", light=True)
    _ftr(sl)

    NW, NH = 1.55, 0.78
    FLOW_Y = 2.7

    steps = [
        (0.25,  FLOW_Y, "Dev / Squad",    "VS Code + Copilot",    AZ_BLUE_LT,  AZ_BLUE,   1),
        (2.15,  FLOW_Y, "Git Repo ADO",   "Azure DevOps / Git",   AZ_BLUE_LT,  AZ_BLUE,   2),
        (4.05,  FLOW_Y, "CI Pipeline",    "Build + Tests + Lint",  AZ_RED_LT,   AZ_RED,    3),
        (5.95,  FLOW_Y, "ACR",            "Docker Push / Tag",     AZ_ORG_LT,   AZ_ORANGE, 4),
        (7.85,  FLOW_Y, "Container Apps", "ACA Deploy Azure",      AZ_PURPLE_LT,AZ_PURPLE, 5),
        (9.75,  FLOW_Y, "AI Foundry",     "Agent Registration",    AZ_PURPLE_LT,AZ_PURPLE, 6),
        (11.65, FLOW_Y, "Teams Live",     "Nova versao ativa",     AZ_GREEN_LT, AZ_GREEN,  7),
    ]

    for x, y, t, s, f, b, n in steps:
        shp = _rect(sl, x, y, NW, NH, f, b, 2.0, rx=True)
        _label_in(shp, t, s, size1=8.5, size2=6.5, bold=True, col1=b)
        _badge(sl, x + NW - 0.18, y - 0.17, n)

    # main flow arrows
    for i in range(len(steps) - 1):
        x1 = steps[i][0] + NW
        x2 = steps[i + 1][0]
        _arrow(sl, x1, FLOW_Y + NH / 2, x2, FLOW_Y + NH / 2, GRAY_DARK, pt=1.5)

    arr_labels = ["PR / Commit", "Build Image", "Push", "Deploy ACA", "Register", "Activate"]
    for i, lbl in enumerate(arr_labels):
        mx = (steps[i][0] + NW + steps[i + 1][0]) / 2
        _label(sl, lbl, mx - 0.5, FLOW_Y + NH + 0.05, 1.0, 0.22,
               size=6, italic=True, color=GRAY_DARK, align=PP_ALIGN.CENTER)

    # Key Vault branch (above)
    _kv_icon(sl, 4.52, 0.68, "Key Vault", step_n="S")
    _label(sl, "Inject secrets during\nCI/CD pipeline run", 4.3, 1.85, 1.9, 0.32,
           size=6.5, italic=True, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    _arrow(sl, 5.07, 1.78, 5.07, FLOW_Y, GOLD_DK, pt=1.5)

    # Observability section
    _label(sl, "Observabilidade em tempo real", 3.5, 4.72, 5.0, 0.26,
           size=9, bold=True, color=AZ_TEAL, align=PP_ALIGN.CENTER)

    obs_items = [
        (3.2,  "Azure Monitor\n+ App Insights",  AZ_TEAL,   AZ_TEAL_LT),
        (5.55, "Defender\nfor Cloud",             AZ_TEAL,   AZ_TEAL_LT),
        (7.9,  "DORA Metrics\n+ SLO Alerts",      AZ_TEAL,   AZ_TEAL_LT),
    ]
    for ox, ot, oc, obg in obs_items:
        s = _rect(sl, ox, 5.05, 2.1, 0.75, obg, oc, 1.5, rx=True)
        _label_in(s, ot, size1=8.5, bold=True, col1=oc)
        _arrow(sl, ox + 1.05, FLOW_Y + NH, ox + 1.05, 5.05, oc, pt=1.0)

    # Cost summary boxes
    for i, (cl, cv, cc) in enumerate([
        ("MVP (10-20 users)",       "~$955 / mes",   AZ_GREEN),
        ("Producao (50-100 users)", "~$4.655 / mes", AZ_PURPLE),
        ("Prazo de Implementacao",  "3-5 meses",     AZ_TEAL),
    ]):
        bx = 0.25 + i * 3.05
        _rect(sl, bx, 6.05, 2.78, 0.85, WHITE, cc, 1.5, rx=True)
        _label(sl, cl, bx + 0.1, 6.1, 2.58, 0.26, size=7.5, color=GRAY_DARK)
        _label(sl, cv, bx + 0.1, 6.38, 2.58, 0.45, size=16, bold=True, color=cc)

    _label(sl, "Microsoft Azure", 0.12, 6.88, 1.55, 0.28,
           size=7, bold=True, color=AZ_BLUE)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — ROADMAP
# ═══════════════════════════════════════════════════════════════════════════════

def slide_roadmap(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bk = sl.background
    bk.fill.solid()
    bk.fill.fore_color.rgb = WHITE

    _hdr(sl, "Roadmap de Implementacao  --  4 Fases / 3 a 5 Meses",
         "Sequencia de atividades para disponibilizacao completa no Azure + Teams", light=True)
    _ftr(sl, "Total: 3-5 meses  |  MVP: ~$955/mes  |  Producao: ~$4.655/mes  |  KARE Project -- Maio 2026")

    phases = [
        dict(n=1, t="Fundacao",         p="Sem 1-6",   c=AZ_BLUE,   bg=AZ_BLUE_LT, items=[
            "Provisionar Azure (Sub, RG, RBAC, policies)",
            "Setup Azure OpenAI + Key Vault + secrets",
            "Containerizar Context Engine FastAPI",
            "Configurar AI Foundry Service",
            "Migrar SQLite para Cosmos DB",
            "Setup Container Registry + Azure DevOps",
        ]),
        dict(n=2, t="Integracao Core",  p="Sem 7-12",  c=AZ_PURPLE, bg=AZ_PURPLE_LT, items=[
            "Registrar 34 agentes no AI Foundry",
            "Containerizar os 3 MCP Servers",
            "Implementar Azure AI Bot Service",
            "Migrar busca para Azure AI Search",
            "Setup API Management + Entra ID SSO",
            "Testes de integracao end-to-end",
        ]),
        dict(n=3, t="Teams + Producao", p="Sem 13-17", c=AZ_TEAL,   bg=AZ_TEAL_LT, items=[
            "Criar Teams App KARE (Bot + Tab)",
            "Publicar bot no canal Teams (SSO corporativo)",
            "47 slash-commands ativos no Teams",
            "Setup Front Door + WAF + DDoS",
            "Pipeline CI/CD no Azure DevOps",
            "Testes de carga + ajuste PTU OpenAI",
        ]),
        dict(n=4, t="Observabilidade",  p="Sem 18-19", c=AZ_GREEN,  bg=AZ_GREEN_LT, items=[
            "Dashboard Azure Monitor + App Insights",
            "Alertas SLO e runbooks operacionais",
            "Auditoria Defender for Cloud",
            "Documentacao e onboarding dos squads",
            "Go-live gradual (beta  ->  producao)",
        ]),
    ]

    PW = (13.333 - 0.35 - 0.14 * 3) / 4

    # Timeline bar
    _rect(sl, 0.18, 0.62, 13.0, 0.14, GRAY_LIGHT)
    for i, (ml, mc) in enumerate(zip(
        ["Fundacao", "Core", "Teams", "Go-Live"],
        [AZ_BLUE, AZ_PURPLE, AZ_TEAL, AZ_GREEN]
    )):
        mx = 0.18 + (i + 1) * (13.0 / 4)
        _oval(sl, mx - 0.14, 0.55, 0.28, 0.28, mc)
        _label(sl, ml, mx - 0.45, 0.85, 0.9, 0.2,
               size=6.5, bold=True, color=mc, align=PP_ALIGN.CENTER)

    for pi, ph in enumerate(phases):
        px = 0.18 + pi * (PW + 0.14)
        py = 1.1

        _rect(sl, px, py, PW, 1.18, ph["c"])
        n_oval = _oval(sl, px + PW / 2 - 0.23, py + 0.08, 0.46, 0.46, WHITE)
        _label(sl, str(ph["n"]), px + PW / 2 - 0.23, py + 0.08, 0.46, 0.46,
               size=14, bold=True, color=ph["c"], align=PP_ALIGN.CENTER)
        _label(sl, ph["t"], px + 0.05, py + 0.58, PW - 0.1, 0.3,
               size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _label(sl, ph["p"], px + 0.05, py + 0.89, PW - 0.1, 0.22,
               size=8, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

        IH = 5.12
        _rect(sl, px, py + 1.18, PW, IH, ph["bg"], ph["c"], 0.75)

        rh = IH / len(ph["items"])
        for ii, item in enumerate(ph["items"]):
            iy = py + 1.18 + ii * rh
            _oval(sl, px + 0.1, iy + rh / 2 - 0.07, 0.13, 0.13, ph["c"])
            _label(sl, item, px + 0.27, iy + 0.05, PW - 0.34, rh - 0.1,
                   size=7.5, color=TEXT_DARK, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  BUILD
# ═══════════════════════════════════════════════════════════════════════════════

def build(dest=None):
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_capa(prs)
    slide_fluxo(prs)
    slide_stack(prs)
    slide_devops(prs)
    slide_roadmap(prs)

    if not dest:
        base = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
        out_dir = os.path.join(base, "_outputs", "kare-azure", "upstream")
        os.makedirs(out_dir, exist_ok=True)
        dest = os.path.join(out_dir, "KARE_Azure_Architecture.pptx")

    prs.save(dest)
    return dest


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="PPT Arquitetura Azure KARE v2 -- flow diagram")
    parser.add_argument("--dest", default=None)
    args = parser.parse_args()
    path = build(args.dest)
    print(f"[OK] {path}")
    print("     5 slides: Capa | Fluxo GitOps-IA | Stack | DevOps CI-CD | Roadmap")
