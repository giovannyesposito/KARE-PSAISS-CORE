#!/usr/bin/env python3
"""
Canvas to KARE Context Converter
=================================

Extracts text, tables, and structure from supported document formats and emits
KARE-ready Markdown with canvas-type heuristics and frontmatter.

Supported formats:
  .pptx / .ppt  — PowerPoint (python-pptx; .ppt requires Win32 + PowerPoint)
  .docx / .doc  — Word document (python-docx; .doc requires Win32 + Word)
  .pdf          — PDF document (pypdf)
  .xlsx         — Excel spreadsheet (openpyxl)

Usage:
    python .agent/scripts/generators/ppt_to_kare.py input.pptx
    python .agent/scripts/generators/ppt_to_kare.py canvas.pptx --output uploads/CANVAS.md --kare-context
    python .agent/scripts/generators/ppt_to_kare.py brief.docx  --output uploads/BRIEF.md  --kare-context
    python .agent/scripts/generators/ppt_to_kare.py report.pdf  --output uploads/REPORT.md --kare-context
    python .agent/scripts/generators/ppt_to_kare.py matrix.xlsx --output uploads/MATRIX.md --kare-context
    python .agent/scripts/generators/ppt_to_kare.py legacy.ppt  --output uploads/LEGACY.md --include-notes
"""

from __future__ import annotations

import argparse
import json
import os
import re
import sys
import tempfile
import unicodedata
from dataclasses import asdict, dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation

PPTX_SAVE_AS = 24

INITIATIVE_KEYWORDS = {
    "iniciativa",
    "programa",
    "problema",
    "objetivo",
    "beneficio",
    "beneficios",
    "escopo",
    "fora de escopo",
    "stakeholder",
    "stakeholders",
    "risco",
    "riscos",
    "solicitante",
    "diretoria",
    "impacto no negocio",
    "impacto",
    "kpi",
    "success metrics",
}

FEATURE_KEYWORDS = {
    "feature",
    "funcionalidade",
    "user story",
    "story",
    "stories",
    "epico",
    "epic",
    "criterio de aceite",
    "criterios de aceite",
    "acceptance criteria",
    "gherkin",
    "cenario",
    "cenario de teste",
    "fluxo da feature",
    "regra funcional",
    "rf",
    "tela",
    "backend",
    "frontend",
}

CANVAS_PREFIX = {
    "initiative": "INITIATIVE",
    "feature": "FEATURE",
    "unknown": "PRESENTATION",
}

CANVAS_NAME_HINTS = {
    "initiative": {"INITIATIVE", "INICIATIVA", "CANVAS_DE_INICIATIVA", "CANVAS_INICIATIVA"},
    "feature": {"FEATURE", "FUNCIONALIDADE", "CANVAS_DE_FEATURE", "CANVAS_FEATURE"},
}

SUPPORTED_EXTENSIONS: frozenset[str] = frozenset({".ppt", ".pptx", ".doc", ".docx", ".pdf", ".xlsx"})

SOURCE_INFO: dict[str, dict[str, str]] = {
    "pptx": {
        "unit": "Slide",
        "count_label": "Quantidade de slides",
        "doc_desc": "apresentação PowerPoint",
    },
    "docx": {
        "unit": "Seção",
        "count_label": "Quantidade de seções",
        "doc_desc": "documento Word",
    },
    "pdf": {
        "unit": "Página",
        "count_label": "Quantidade de páginas",
        "doc_desc": "documento PDF",
    },
    "xlsx": {
        "unit": "Planilha",
        "count_label": "Quantidade de planilhas",
        "doc_desc": "planilha Excel",
    },
}
_DEFAULT_SOURCE_INFO: dict[str, str] = {
    "unit": "Seção",
    "count_label": "Quantidade de seções",
    "doc_desc": "documento",
}


@dataclass
class SlideContent:
    number: int
    title: str
    text_blocks: list[str] = field(default_factory=list)
    tables: list[list[list[str]]] = field(default_factory=list)
    notes: list[str] = field(default_factory=list)


@dataclass
class PresentationContent:
    source_file: str
    extracted_at: str
    slide_count: int
    title: str
    canvas_type: str
    heuristic_confidence: float
    heuristic_matches: dict[str, list[str]]
    slides: list[SlideContent]
    source_type: str = "pptx"


@dataclass
class CanvasClassification:
    canvas_type: str
    confidence: float
    matches: dict[str, list[str]]


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Convert canvas documents (.pptx, .docx, .pdf, .xlsx) into KARE-friendly Markdown.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=(
            "Examples:\n"
            "  python .agent/scripts/generators/ppt_to_kare.py canvas.pptx\n"
            "  python .agent/scripts/generators/ppt_to_kare.py canvas.pptx --output uploads/CANVAS.md --kare-context\n"
            "  python .agent/scripts/generators/ppt_to_kare.py brief.docx  --output uploads/BRIEF.md  --kare-context\n"
            "  python .agent/scripts/generators/ppt_to_kare.py report.pdf  --output uploads/REPORT.md --kare-context\n"
            "  python .agent/scripts/generators/ppt_to_kare.py matrix.xlsx --output uploads/MATRIX.md --kare-context\n"
            "  python .agent/scripts/generators/ppt_to_kare.py legacy.ppt  --format json --include-notes"
        ),
    )
    parser.add_argument(
        "input",
        help="Path to input file (.pptx, .ppt, .docx, .doc, .pdf, or .xlsx)",
    )
    parser.add_argument(
        "--output",
        help="Optional output file. If omitted, content is written to stdout.",
    )
    parser.add_argument(
        "--format",
        choices=("markdown", "text", "json"),
        default="markdown",
        help="Output format. Default: markdown",
    )
    parser.add_argument(
        "--include-notes",
        action="store_true",
        help="Include speaker notes when available.",
    )
    parser.add_argument(
        "--kare-context",
        action="store_true",
        help="Emit Markdown with KARE context frontmatter.",
    )
    parser.add_argument(
        "--title",
        help="Override the document title used in the generated output.",
    )
    parser.add_argument(
        "--save-to-context",
        action="store_true",
        help="Save the generated markdown automatically into the workspace uploads/ folder.",
    )
    parser.add_argument(
        "--context-name",
        help="Optional KARE context file name. It will be normalized to UPPER_SNAKE_CASE.",
    )
    return parser


def ensure_supported_extension(input_path: Path) -> None:
    if input_path.suffix.lower() not in SUPPORTED_EXTENSIONS:
        supported_str = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValueError(f"Unsupported file type '{input_path.suffix}'. Supported: {supported_str}")


def convert_ppt_to_pptx(input_path: Path) -> Path:
    if os.name != "nt":
        raise RuntimeError("Legacy .ppt conversion is only supported on Windows.")

    try:
        import win32com.client  # type: ignore[import-not-found]
    except ImportError as exc:
        raise RuntimeError(
            "Reading .ppt requires pywin32 and Microsoft PowerPoint installed. "
            "Install pywin32 or convert the file to .pptx first."
        ) from exc

    powerpoint = None
    presentation = None
    temp_dir = Path(tempfile.mkdtemp(prefix="kare-ppt-"))
    output_path = temp_dir / f"{input_path.stem}.pptx"

    try:
        powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
        powerpoint.Visible = 0
        presentation = powerpoint.Presentations.Open(str(input_path), WithWindow=False)
        presentation.SaveAs(str(output_path), PPTX_SAVE_AS)
    except Exception as exc:
        raise RuntimeError(
            "Failed to convert .ppt to .pptx using Microsoft PowerPoint."
        ) from exc
    finally:
        if presentation is not None:
            presentation.Close()
        if powerpoint is not None:
            powerpoint.Quit()

    return output_path


def normalize_text(value: str) -> str:
    lines = [line.strip() for line in value.splitlines()]
    cleaned_lines = [line for line in lines if line]
    return "\n".join(cleaned_lines).strip()


def extract_table(shape: Any) -> list[list[str]]:
    rows: list[list[str]] = []
    for row in shape.table.rows:
        rows.append([normalize_text(cell.text) for cell in row.cells])
    return rows


def extract_notes(slide: Any) -> list[str]:
    notes: list[str] = []
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return notes

    for shape in notes_slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = normalize_text(shape.text)
        if text:
            notes.append(text)
    return deduplicate(notes)


def deduplicate(values: list[str]) -> list[str]:
    unique_values: list[str] = []
    seen: set[str] = set()
    for value in values:
        if value in seen:
            continue
        seen.add(value)
        unique_values.append(value)
    return unique_values


def normalize_for_matching(value: str) -> str:
    normalized = unicodedata.normalize("NFKD", value)
    ascii_only = normalized.encode("ascii", "ignore").decode("ascii")
    return re.sub(r"\s+", " ", ascii_only.lower()).strip()


def collect_presentation_text(slides: list[SlideContent]) -> str:
    parts: list[str] = []
    for slide in slides:
        parts.append(slide.title)
        parts.extend(slide.text_blocks)
        parts.extend(slide.notes)
        for table in slide.tables:
            for row in table:
                parts.extend(row)
    return "\n".join(parts)


def find_keyword_matches(content: str, keywords: set[str]) -> list[str]:
    matches = [keyword for keyword in sorted(keywords) if keyword in content]
    return matches


def classify_canvas_type(slides: list[SlideContent], source_file: Path) -> CanvasClassification:
    combined_text = normalize_for_matching(collect_presentation_text(slides))
    source_hint = normalize_for_matching(source_file.stem)

    initiative_matches = find_keyword_matches(combined_text, INITIATIVE_KEYWORDS)
    feature_matches = find_keyword_matches(combined_text, FEATURE_KEYWORDS)

    if "iniciativa" in source_hint:
        initiative_matches.append("filename:iniciativa")
    if "feature" in source_hint or "funcionalidade" in source_hint:
        feature_matches.append("filename:feature")

    initiative_score = len(initiative_matches)
    feature_score = len(feature_matches)

    if initiative_score == 0 and feature_score == 0:
        return CanvasClassification(
            canvas_type="unknown",
            confidence=0.0,
            matches={"initiative": [], "feature": []},
        )

    if initiative_score > feature_score:
        total = initiative_score + feature_score
        return CanvasClassification(
            canvas_type="initiative",
            confidence=initiative_score / total,
            matches={"initiative": deduplicate(initiative_matches), "feature": deduplicate(feature_matches)},
        )

    if feature_score > initiative_score:
        total = initiative_score + feature_score
        return CanvasClassification(
            canvas_type="feature",
            confidence=feature_score / total,
            matches={"initiative": deduplicate(initiative_matches), "feature": deduplicate(feature_matches)},
        )

    return CanvasClassification(
        canvas_type="unknown",
        confidence=0.5,
        matches={"initiative": deduplicate(initiative_matches), "feature": deduplicate(feature_matches)},
    )


def extract_slide(slide: Any, slide_number: int, include_notes: bool) -> SlideContent:
    title = ""
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None:
        title = normalize_text(title_shape.text)

    text_blocks: list[str] = []
    tables: list[list[list[str]]] = []

    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = normalize_text(shape.text)
            if text and text != title:
                text_blocks.append(text)

        if getattr(shape, "has_table", False):
            table = extract_table(shape)
            if table:
                tables.append(table)

    notes = extract_notes(slide) if include_notes else []

    return SlideContent(
        number=slide_number,
        title=title or f"Slide {slide_number}",
        text_blocks=deduplicate(text_blocks),
        tables=tables,
        notes=notes,
    )


def _parse_pptx_content(
    input_path: Path,
    include_notes: bool,
) -> tuple[list[SlideContent], str]:
    """Parse .pptx (and .ppt via COM) using python-pptx."""
    working_path = input_path
    cleanup_path: Path | None = None

    if input_path.suffix.lower() == ".ppt":
        working_path = convert_ppt_to_pptx(input_path)
        cleanup_path = working_path

    prs_title = ""
    slides: list[SlideContent] = []
    try:
        presentation = Presentation(str(working_path))
        slides = [
            extract_slide(slide, index, include_notes)
            for index, slide in enumerate(presentation.slides, start=1)
        ]
        try:
            prs_title = (presentation.core_properties.title or "").strip()
        except Exception:
            pass
    finally:
        if cleanup_path is not None:
            cleanup_path.unlink(missing_ok=True)
            cleanup_path.parent.rmdir()

    return slides, prs_title


def _convert_doc_to_docx(input_path: Path) -> Path:
    """Convert legacy .doc to .docx using Microsoft Word on Windows."""
    if os.name != "nt":
        raise RuntimeError("Legacy .doc conversion is only supported on Windows.")
    try:
        import win32com.client as _wc  # type: ignore[import-not-found]
    except ImportError as exc:
        raise RuntimeError(
            "Reading .doc requires pywin32 and Microsoft Word installed. "
            "Install pywin32 or convert the file to .docx first."
        ) from exc

    DOCX_SAVE_FORMAT = 16  # wdFormatXMLDocument
    word = None
    document = None
    temp_dir = Path(tempfile.mkdtemp(prefix="kare-doc-"))
    output_path = temp_dir / f"{input_path.stem}.docx"
    try:
        word = _wc.DispatchEx("Word.Application")
        word.Visible = False
        document = word.Documents.Open(str(input_path), ReadOnly=True)
        document.SaveAs2(str(output_path), FileFormat=DOCX_SAVE_FORMAT)
    except Exception as exc:
        raise RuntimeError("Failed to convert .doc to .docx using Microsoft Word.") from exc
    finally:
        if document is not None:
            document.Close(False)
        if word is not None:
            word.Quit()
    return output_path


def _parse_docx_content(
    input_path: Path,
    include_notes: bool,
) -> tuple[list[SlideContent], str]:
    """Parse .docx/.doc files using python-docx."""
    try:
        from docx import Document  # type: ignore[import-not-found]
        from docx.oxml.ns import qn  # type: ignore[import-not-found]
        from docx.text.paragraph import Paragraph as _DocParagraph  # type: ignore[import-not-found]
        from docx.table import Table as _DocTable  # type: ignore[import-not-found]
    except ImportError as exc:
        raise RuntimeError(
            "Reading .docx requires python-docx. Install with: pip install python-docx"
        ) from exc

    temp_cleanup: Path | None = None
    actual_path = input_path

    if input_path.suffix.lower() == ".doc":
        actual_path = _convert_doc_to_docx(input_path)
        temp_cleanup = actual_path

    try:
        doc = Document(str(actual_path))

        def _iter_blocks(parent: Any) -> Any:
            from docx.table import _Cell  # type: ignore[import-not-found]
            parent_elm = parent._tc if isinstance(parent, _Cell) else parent.element.body
            for child in parent_elm.iterchildren():
                if child.tag == qn("w:p"):
                    yield _DocParagraph(child, parent)
                elif child.tag == qn("w:tbl"):
                    yield _DocTable(child, parent)

        HEADING_PREFIXES = ("heading", "título", "titulo", "title")
        sections: list[SlideContent] = []
        current_title = ""
        current_blocks: list[str] = []
        current_tables: list[list[list[str]]] = []
        section_number = 0

        def _flush_section() -> None:
            nonlocal current_title, current_blocks, current_tables, section_number
            if current_title or current_blocks or current_tables:
                section_number += 1
                sections.append(
                    SlideContent(
                        number=section_number,
                        title=current_title or f"Seção {section_number}",
                        text_blocks=deduplicate([b for b in current_blocks if b]),
                        tables=list(current_tables),
                        notes=[],
                    )
                )
            current_title = ""
            current_blocks = []
            current_tables = []

        for item in _iter_blocks(doc):
            if isinstance(item, _DocParagraph):
                text = normalize_text(item.text)
                if not text:
                    continue
                style_name = (item.style.name or "").lower() if item.style else ""
                if any(style_name.startswith(p) for p in HEADING_PREFIXES):
                    _flush_section()
                    current_title = text
                else:
                    current_blocks.append(text)
            elif isinstance(item, _DocTable):
                rows: list[list[str]] = [
                    [normalize_text(cell.text) for cell in row.cells]
                    for row in item.rows
                ]
                if rows:
                    current_tables.append(rows)

        _flush_section()

        meta_title = ""
        try:
            meta_title = (doc.core_properties.title or "").strip()
        except Exception:
            pass

    finally:
        if temp_cleanup is not None:
            temp_cleanup.unlink(missing_ok=True)
            try:
                temp_cleanup.parent.rmdir()
            except Exception:
                pass

    return sections, meta_title


def _parse_pdf_content(
    input_path: Path,
    include_notes: bool,
) -> tuple[list[SlideContent], str]:
    """Parse PDF files using pypdf."""
    try:
        from pypdf import PdfReader  # type: ignore[import-not-found]
    except ImportError as exc:
        raise RuntimeError(
            "Reading PDF requires pypdf. Install with: pip install pypdf"
        ) from exc

    reader = PdfReader(str(input_path))
    pages: list[SlideContent] = []

    for page_number, page in enumerate(reader.pages, start=1):
        raw_text = page.extract_text() or ""
        raw_blocks = [normalize_text(b) for b in re.split(r"\n{2,}", raw_text)]
        text_blocks = [b for b in raw_blocks if b]
        if not text_blocks:
            continue

        first = text_blocks[0]
        if len(first) <= 120 and "\n" not in first:
            title = first
            body = text_blocks[1:]
        else:
            title = f"Página {page_number}"
            body = text_blocks

        pages.append(
            SlideContent(
                number=page_number,
                title=title,
                text_blocks=body,
                tables=[],
                notes=[],
            )
        )

    meta_title = ""
    try:
        meta = reader.metadata
        if meta:
            meta_title = (meta.title or "").strip()
    except Exception:
        pass

    return pages, meta_title


def _parse_xlsx_content(
    input_path: Path,
    include_notes: bool,
) -> tuple[list[SlideContent], str]:
    """Parse .xlsx files using openpyxl."""
    try:
        import openpyxl  # type: ignore[import-not-found]
    except ImportError as exc:
        raise RuntimeError(
            "Reading .xlsx requires openpyxl. Install with: pip install openpyxl"
        ) from exc

    workbook = openpyxl.load_workbook(str(input_path), data_only=True)
    sheets: list[SlideContent] = []

    for sheet_number, sheet_name in enumerate(workbook.sheetnames, start=1):
        ws = workbook[sheet_name]
        table_rows: list[list[str]] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                table_rows.append(cells)

        if not table_rows:
            continue

        # Extract header row + first-column labels as text for heuristic classification
        header_text = " ".join(c for c in table_rows[0] if c.strip())
        first_col = [r[0] for r in table_rows[1:] if r and r[0].strip()]
        text_blocks: list[str] = []
        if header_text:
            text_blocks.append(header_text)
        if first_col:
            text_blocks.append(" ".join(first_col[:20]))

        sheets.append(
            SlideContent(
                number=sheet_number,
                title=sheet_name,
                text_blocks=text_blocks,
                tables=[table_rows],
                notes=[],
            )
        )

    meta_title = ""
    try:
        if workbook.properties and workbook.properties.title:
            meta_title = workbook.properties.title.strip()
    except Exception:
        pass

    workbook.close()
    return sheets, meta_title


def parse_presentation(input_path: Path, include_notes: bool, title_override: str | None) -> PresentationContent:
    ext = input_path.suffix.lower()

    if ext in {".ppt", ".pptx"}:
        slides, doc_title = _parse_pptx_content(input_path, include_notes)
        source_type = "pptx"
    elif ext in {".doc", ".docx"}:
        slides, doc_title = _parse_docx_content(input_path, include_notes)
        source_type = "docx"
    elif ext == ".pdf":
        slides, doc_title = _parse_pdf_content(input_path, include_notes)
        source_type = "pdf"
    elif ext == ".xlsx":
        slides, doc_title = _parse_xlsx_content(input_path, include_notes)
        source_type = "xlsx"
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    inferred_title = title_override or find_title(slides) or doc_title or input_path.stem.replace("_", " ")
    classification = classify_canvas_type(slides, input_path)

    return PresentationContent(
        source_file=str(input_path),
        extracted_at=datetime.now().isoformat(timespec="seconds"),
        slide_count=len(slides),
        title=inferred_title,
        canvas_type=classification.canvas_type,
        heuristic_confidence=round(classification.confidence, 3),
        heuristic_matches=classification.matches,
        slides=slides,
        source_type=source_type,
    )


def find_title(slides: list[SlideContent]) -> str:
    _default_prefixes = ("Slide ", "Seção ", "Página ", "Planilha ")
    for slide in slides:
        if slide.title and not any(slide.title.startswith(p) for p in _default_prefixes):
            return slide.title
    return ""


def to_upper_snake(value: str) -> str:
    normalized = unicodedata.normalize("NFKD", value)
    ascii_only = normalized.encode("ascii", "ignore").decode("ascii")
    collapsed = re.sub(r"[^A-Za-z0-9]+", "_", ascii_only).strip("_")
    return collapsed.upper() or "PRESENTATION_IMPORT"


def find_workspace_root() -> Path:
    search_roots = [Path.cwd().resolve(), Path(__file__).resolve().parents[2]]
    for start in search_roots:
        for candidate in [start, *start.parents]:
            if (candidate / ".agent").exists() and (candidate / "context").exists():
                return candidate
    return Path.cwd().resolve()


def resolve_context_output_path(content: PresentationContent, context_name: str | None) -> Path:
    workspace_root = find_workspace_root()
    context_dir = workspace_root / "context"
    if context_name:
        file_stem = to_upper_snake(context_name)
    else:
        prefix = CANVAS_PREFIX.get(content.canvas_type, "PRESENTATION")
        base_name = to_upper_snake(content.title or Path(content.source_file).stem)
        name_hints = CANVAS_NAME_HINTS.get(content.canvas_type, set())
        has_equivalent_prefix = base_name.startswith(f"{prefix}_") or any(
            base_name.startswith(f"{hint}_") or base_name == hint for hint in name_hints
        )
        if has_equivalent_prefix:
            file_stem = base_name
        else:
            file_stem = f"{prefix}_{base_name}"
    return context_dir / f"{file_stem}.md"


def render_table_markdown(table: list[list[str]]) -> list[str]:
    if not table:
        return []

    header = table[0]
    separator = ["---"] * len(header)
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    for row in table[1:]:
        padded_row = row + [""] * (len(header) - len(row))
        lines.append("| " + " | ".join(padded_row[: len(header)]) + " |")
    return lines


def render_markdown(content: PresentationContent, kare_context: bool) -> str:
    _src = SOURCE_INFO.get(content.source_type, _DEFAULT_SOURCE_INFO)
    lines: list[str] = []
    if kare_context:
        lines.extend(
            [
                "---",
                "type: context-document",
                f'last_updated: "{datetime.now().date().isoformat()}"',
                'author: "canvas_to_kare.py"',
                'version: "1.0"',
                f'canvas_type: "{content.canvas_type}"',
                f'heuristic_confidence: {content.heuristic_confidence}',
                "---",
                "",
            ]
        )

    lines.extend(
        [
            f"# {content.title}",
            "",
            f"> Documento gerado automaticamente a partir de {_src['doc_desc']} para uso no KARE.",
            "",
            "## Metadados",
            "",
            f"- Arquivo de origem: {content.source_file}",
            f"- Extraído em: {content.extracted_at}",
            f"- {_src['count_label']}: {content.slide_count}",
            f"- Tipo de canvas: {content.canvas_type}",
            f"- Confiança da heurística: {content.heuristic_confidence}",
            "",
        ]
    )

    if content.heuristic_matches["initiative"] or content.heuristic_matches["feature"]:
        lines.extend(
            [
                "## Detecção Heurística",
                "",
                f"- Correspondências de initiative: {', '.join(content.heuristic_matches['initiative']) or 'nenhuma'}",
                f"- Correspondências de feature: {', '.join(content.heuristic_matches['feature']) or 'nenhuma'}",
                "",
            ]
        )

    for slide in content.slides:
        lines.append(f"## {_src['unit']} {slide.number} — {slide.title}")
        lines.append("")

        if slide.text_blocks:
            lines.append("### Texto")
            lines.append("")
            for block in slide.text_blocks:
                lines.append(block)
                lines.append("")

        if slide.tables:
            lines.append("### Tabelas")
            lines.append("")
            for index, table in enumerate(slide.tables, start=1):
                lines.append(f"Tabela {index}")
                lines.extend(render_table_markdown(table))
                lines.append("")

        if slide.notes:
            lines.append("### Notas do Apresentador")
            lines.append("")
            for note in slide.notes:
                lines.append(note)
                lines.append("")

    return "\n".join(lines).strip() + "\n"


def render_text(content: PresentationContent) -> str:
    _src = SOURCE_INFO.get(content.source_type, _DEFAULT_SOURCE_INFO)
    lines = [
        content.title,
        "=" * len(content.title),
        f"Arquivo de origem: {content.source_file}",
        f"Extraído em: {content.extracted_at}",
        f"{_src['count_label']}: {content.slide_count}",
        f"Tipo de canvas: {content.canvas_type}",
        f"Confiança da heurística: {content.heuristic_confidence}",
        "",
    ]

    for slide in content.slides:
        lines.append(f"{_src['unit']} {slide.number}: {slide.title}")
        lines.append("-" * (len(slide.title) + 10))
        lines.extend(slide.text_blocks)
        if slide.tables:
            for index, table in enumerate(slide.tables, start=1):
                lines.append(f"[Tabela {index}]")
                for row in table:
                    lines.append(" | ".join(row))
        if slide.notes:
            lines.append("[Notas do Apresentador]")
            lines.extend(slide.notes)
        lines.append("")

    return "\n".join(lines).strip() + "\n"


def render_json(content: PresentationContent) -> str:
    return json.dumps(asdict(content), ensure_ascii=False, indent=2) + "\n"


def render_output(content: PresentationContent, output_format: str, kare_context: bool) -> str:
    if output_format == "json":
        return render_json(content)
    if output_format == "text":
        return render_text(content)
    return render_markdown(content, kare_context)


def write_output(rendered_output: str, output_path: Path | None) -> Path | None:
    if output_path is None:
        sys.stdout.write(rendered_output)
        return None

    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(rendered_output, encoding="utf-8")
    return output_path


def main() -> int:
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8")
    if hasattr(sys.stderr, "reconfigure"):
        sys.stderr.reconfigure(encoding="utf-8")

    parser = build_parser()
    args = parser.parse_args()

    if args.save_to_context and args.output:
        parser.error("Use either --output or --save-to-context, not both.")

    if args.save_to_context and args.format != "markdown":
        parser.error("--save-to-context only supports markdown output.")

    input_path = Path(args.input).expanduser().resolve()
    if not input_path.exists():
        parser.error(f"Input file does not exist: {input_path}")

    try:
        ensure_supported_extension(input_path)
        content = parse_presentation(input_path, args.include_notes, args.title)
        use_kare_context = args.kare_context or args.save_to_context
        rendered_output = render_output(content, args.format, use_kare_context)
        output_path = None
        if args.output:
            output_path = Path(args.output).expanduser().resolve()
        elif args.save_to_context:
            output_path = resolve_context_output_path(content, args.context_name)

        written_path = write_output(rendered_output, output_path)
        if written_path is not None:
            sys.stdout.write(
                f"[OK] Generated: {written_path} (canvas_type={content.canvas_type}, confidence={content.heuristic_confidence})\n"
            )
    except Exception as exc:
        sys.stderr.write(f"Error: {exc}\n")
        return 1

    return 0


if __name__ == "__main__":
    raise SystemExit(main())