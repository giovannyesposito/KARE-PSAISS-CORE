#!/usr/bin/env python3
"""
kare_contexto.py � Ingest�o de arquivos externos no KARE Context Engine (RAG)

Uso:
  python kare_contexto.py --file <path>  --context <slug> [op��es]
  python kare_contexto.py --dir  <path>  --context <slug> [--recursive]

Formatos suportados:
  Texto     ? .md  .txt  .rst
  Dados     ? .csv  .json
  Office    ? .docx  .pptx  .xlsx  .xls
  PDF       ? .pdf

Depend�ncias opcionais (instale conforme necess�rio):
  pip install pdfplumber python-docx python-pptx openpyxl
"""

import sys
import os
import json
import argparse
import urllib.request
import urllib.parse
import urllib.error
from pathlib import Path
from typing import Optional
import re
import sys

sys.path.insert(0, str(Path(__file__).parent))
from rag_auth import get_auth_headers as _rag_auth  # noqa: E402

# ---------------------------------------------
API_BASE = "http://localhost:8000"

SUPPORTED_EXTS = {
    ".md", ".txt", ".rst",
    ".csv", ".json",
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx", ".xls",
}

NODE_TYPES = ["auto", "artifact", "symbol", "decision", "concept", "context"]

# ---------------------------------------------
# Auto-classifier
# ---------------------------------------------

# Keyword sets (lowercase) for each node type
_DECISION_KW = {
    "adr", "adr-", "decis�o", "decisao", "decidimos", "optamos por",
    "tradeoff", "trade-off", "status: accepted", "status: proposed",
    "## decision", "## context", "## consequences", "chosen solution",
    "rejected alternatives", "alternativas rejeitadas",
}
_SYMBOL_KW = {
    "gloss�rio", "glossario", "glossary", "sigla", "terminologia",
    "entidade", "defini��o:", "definicao:", "| sigla |", "| term |",
    "| termo |", "| symbol |", "l�xico", "lexico", "dicion�rio", "dicionario",
}
_CONCEPT_KW = {
    "conceito", "overview", "introdu��o", "introducao", "fundamentos",
    "como funciona", "what is", "o que �", "vis�o geral", "visao geral",
}
_CONTEXT_KW = {
    "iniciativa", "programa", "squad", "time de", "project brief",
    "project charter", "kickoff", "kick-off", "kick off",
}


def auto_classify(path: Path, sample: str = "") -> str:
    """Infer best node type from filename and content sample."""
    name_lower = path.stem.lower()
    text = (name_lower + " " + sample[:4000]).lower()

    # Score each type
    scores = {"decision": 0, "symbol": 0, "concept": 0, "context": 0}
    for kw in _DECISION_KW:
        if kw in text:
            scores["decision"] += 1
    for kw in _SYMBOL_KW:
        if kw in text:
            scores["symbol"] += 1
    for kw in _CONCEPT_KW:
        if kw in text:
            scores["concept"] += 1
    for kw in _CONTEXT_KW:
        if kw in text:
            scores["context"] += 1

    best_type = max(scores, key=lambda k: scores[k])
    best_score = scores[best_type]

    # Filename-based overrides (strong signal)
    fn_overrides = [
        (["adr", "decision", "decisao", "decis�o", "decisoes"], "decision"),
        (["glossario", "glossary", "gloss�rio", "siglas", "lexico", "l�xico"], "symbol"),
        (["conceito", "overview", "intro", "fundamentos"], "concept"),
        (["brief", "kickoff", "kick-off", "charter", "iniciativa"], "context"),
    ]
    for keywords, node_type in fn_overrides:
        if any(kw in name_lower for kw in keywords):
            return node_type

    return best_type if best_score > 0 else "artifact"

CHUNK_LABEL = {
    ".pdf":  "P�gs",
    ".pptx": "Slides",
    ".xlsx": "Aba",
    ".xls":  "Aba",
}

# ---------------------------------------------
# HTTP helpers
# ---------------------------------------------

def _api_post(endpoint: str, payload: dict) -> Optional[dict]:
    data = json.dumps(payload, ensure_ascii=False).encode("utf-8")
    req = urllib.request.Request(
        f"{API_BASE}{endpoint}",
        data=data,
        headers={"Content-Type": "application/json", **_rag_auth()},
        method="POST",
    )
    try:
        with urllib.request.urlopen(req, timeout=30) as resp:
            return json.loads(resp.read())
    except urllib.error.HTTPError as e:
        body = e.read().decode("utf-8", errors="replace")
        print(f"  [ERRO HTTP {e.code}] {body[:300]}", file=sys.stderr)
        return None
    except urllib.error.URLError as e:
        print(f"\n[ERRO] API inacess�vel: {e.reason}", file=sys.stderr)
        print(
            "  Verifique os bancos RAG: python .agent/scripts/ai/kare_rag.py status",
            file=sys.stderr,
        )
        sys.exit(1)


def _rebuild_index():
    req = urllib.request.Request(f"{API_BASE}/index/rebuild", method="POST",
                                  headers=_rag_auth())
    try:
        with urllib.request.urlopen(req, timeout=30) as resp:
            r = json.loads(resp.read())
            print(
                f"\n[INDEX] Reconstru�do � {r.get('nodes_indexed','?')} n�s | "
                f"{r.get('duration_ms','?')}ms | {r.get('index_size_mb','?')} MB"
            )
    except Exception as e:
        print(f"\n[AVISO] �ndice n�o reconstru�do: {e}")


def _api_get(path: str):
    """HTTP GET helper � retorna JSON ou None."""
    try:
        req = urllib.request.Request(f"{API_BASE}{path}",
                                      headers={"Accept": "application/json", **_rag_auth()})
        with urllib.request.urlopen(req, timeout=15) as resp:
            return json.loads(resp.read())
    except Exception:
        return None


# Padr�es de ID reconhecidos: INI-001, ADR-025, US-001, FEAT-001, EP-001 etc.
_ID_PATTERN = re.compile(r'\b(INI|ADR|US|FEAT|EP|CAP|EN|TASK|CHG)-\d+\b', re.IGNORECASE)


def auto_link_node(node_id: int, content: str, title: str) -> int:
    """
    Detecta men��es a outros n�s no conte�do deste n� (e vice-versa).
    Cria arestas REFERENCES automaticamente no grafo.
    Retorna o n�mero de arestas criadas.
    """
    all_nodes = _api_get("/nodes?limit=500")
    if not all_nodes:
        return 0

    # Arestas REFERENCES j� existentes (para evitar duplicatas)
    all_edges = _api_get("/edges?limit=2000") or []
    existing = {
        (e["source_id"], e["target_id"])
        for e in all_edges
        if e.get("relation_type") == "references"
    }

    content_lower = content.lower()
    title_lower = title.lower().strip()

    # IDs extra�dos do t�tulo + conte�do deste n� (ex: "INI-001", "ADR-025")
    own_ids = {m.group(0).upper() for m in _ID_PATTERN.finditer(title + " " + content)}

    created = 0

    def _edge(src: int, tgt: int):
        nonlocal created
        pair = (src, tgt)
        if pair in existing:
            return
        r = _api_post("/edges", {
            "source_id": src,
            "target_id": tgt,
            "relation_type": "references",
            "weight": 1.0,
        })
        if r and r.get("id"):
            existing.add(pair)
            created += 1

    for n in all_nodes:
        other_id = n["id"]
        if other_id == node_id:
            continue
        other_title       = (n.get("title") or "").strip()
        other_title_lower = other_title.lower()
        other_content     = (n.get("content") or "")
        other_content_up  = other_content.upper()
        # IDs presentes no t�tulo do outro n�
        other_ids = {m.group(0).upper() for m in _ID_PATTERN.finditer(other_title)}

        # 1. Este ? outro: nosso conte�do menciona o t�tulo do outro n�
        if len(other_title_lower) >= 5 and other_title_lower in content_lower:
            _edge(node_id, other_id)

        # 2. Este ? outro: nosso conte�do menciona um ID presente no t�tulo do outro n�
        for oid in other_ids:
            if oid in content.upper():
                _edge(node_id, other_id)
                break

        # 3. Outro ? este: o conte�do do outro menciona nosso t�tulo
        if len(title_lower) >= 5 and title_lower in other_content.lower():
            _edge(other_id, node_id)

        # 4. Outro ? este: o conte�do do outro menciona um de nossos IDs
        for own_id in own_ids:
            if own_id in other_content_up:
                _edge(other_id, node_id)
                break

    return created


# ---------------------------------------------
# Extractors � retornam list[{"title": str, "content": str}]
# ---------------------------------------------

def _chunked(items: list, size: int, stem: str, label: str) -> list[dict]:
    """Group items (pages / slides) into chunks of `size`."""
    chunks = []
    for i in range(0, len(items), size):
        group = items[i : i + size]
        end   = min(i + size, len(items))
        title = f"{stem} � {label} {i + 1}�{end}"
        chunks.append({"title": title, "content": "\n\n".join(group)})
    return [c for c in chunks if c["content"].strip()]


def extract_txt(path: Path) -> list[dict]:
    text = path.read_text(encoding="utf-8", errors="replace")
    return [{"title": path.stem, "content": text}]


def extract_csv(path: Path) -> list[dict]:
    import csv
    rows = []
    with path.open(encoding="utf-8", errors="replace", newline="") as f:
        reader = csv.reader(f)
        for i, row in enumerate(reader):
            rows.append(" | ".join(str(c) for c in row))
            if i >= 2000:
                rows.append(f"... (truncado em {i} linhas)")
                break
    return [{"title": path.stem, "content": "\n".join(rows)}]


def extract_json(path: Path) -> list[dict]:
    try:
        obj = json.loads(path.read_text(encoding="utf-8", errors="replace"))
        text = json.dumps(obj, ensure_ascii=False, indent=2)
    except Exception:
        text = path.read_text(encoding="utf-8", errors="replace")
    return [{"title": path.stem, "content": text[:60000]}]


def extract_pdf(path: Path) -> list[dict]:
    GROUP = 8
    # Try pdfplumber first
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            pages = [p.extract_text() or "" for p in pdf.pages]
    except ImportError:
        try:
            import PyPDF2
            with open(path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                pages = [p.extract_text() or "" for p in reader.pages]
        except ImportError:
            print(
                "[ERRO] Instale pdfplumber para ler PDFs:\n"
                "  pip install pdfplumber"
            )
            return []
    return _chunked(pages, GROUP, path.stem, "P�gs")


def extract_docx(path: Path) -> list[dict]:
    try:
        from docx import Document
    except ImportError:
        print(
            "[ERRO] Instale python-docx para ler .docx:\n"
            "  pip install python-docx"
        )
        return []

    doc = Document(path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            if para.style.name.startswith("Heading"):
                parts.append(f"\n## {para.text}")
            else:
                parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(c.text.strip() for c in row.cells))
    return [{"title": path.stem, "content": "\n".join(parts)}]


def extract_pptx(path: Path) -> list[dict]:
    GROUP = 5
    try:
        from pptx import Presentation
    except ImportError:
        print(
            "[ERRO] Instale python-pptx para ler .pptx:\n"
            "  pip install python-pptx"
        )
        return []

    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        slides.append(f"--- Slide {i + 1} ---\n" + "\n".join(texts))
    return _chunked(slides, GROUP, path.stem, "Slides")


def extract_xlsx(path: Path) -> list[dict]:
    try:
        import openpyxl
    except ImportError:
        print(
            "[ERRO] Instale openpyxl para ler .xlsx:\n"
            "  pip install openpyxl"
        )
        return []

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    chunks = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        rows = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells))
            if i >= 2000:
                rows.append("... (truncado em 2000 linhas)")
                break
        if rows:
            chunks.append({"title": f"{path.stem} � Aba: {sheet}", "content": "\n".join(rows)})
    wb.close()
    return chunks


# ---------------------------------------------
# Dispatch
# ---------------------------------------------

_EXTRACTORS = {
    ".md":   extract_txt,
    ".txt":  extract_txt,
    ".rst":  extract_txt,
    ".csv":  extract_csv,
    ".json": extract_json,
    ".pdf":  extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".xlsx": extract_xlsx,
    ".xls":  extract_xlsx,
}


def extract_chunks(path: Path) -> list[dict]:
    ext = path.suffix.lower()
    fn  = _EXTRACTORS.get(ext)
    if fn is None:
        print(f"  [PULADO] Extens�o n�o suportada: {ext}")
        return []
    try:
        return fn(path)
    except Exception as e:
        print(f"  [ERRO] Falha ao extrair {path.name}: {e}")
        return []


def resolve_type(node_type: str, path: Path, chunks: list[dict]) -> str:
    """Resolve 'auto' to a concrete type using content from first chunk."""
    if node_type != "auto":
        return node_type
    sample = chunks[0]["content"] if chunks else ""
    resolved = auto_classify(path, sample)
    print(f"  [AUTO] Tipo detectado: {resolved}")
    return resolved


# ---------------------------------------------
# Ingestion
# ---------------------------------------------

def ingest_chunks(
    chunks: list[dict],
    context_slug: str,
    node_type: str,
    title_override: Optional[str] = None,
    source_path: Optional[str] = None,
) -> tuple[int, int, list[dict]]:
    ok = fail = 0
    ingested: list[dict] = []
    for chunk in chunks:
        title   = title_override if (title_override and len(chunks) == 1) else chunk["title"]
        content = chunk["content"]
        if not content.strip():
            print(f"  [VAZIO] {title[:60]}")
            continue

        payload = {
            "type":         node_type,
            "title":        title,
            "content":      content,
            "context_slug": context_slug,
            "metadata": {
                "source":           source_path or title,
                "ingestion_source": "user_context",
                "ingested_via":     "kare_contexto_cli",
            },
        }
        result = _api_post("/ingest", payload)
        if result and result.get("id"):
            print(f"  [OK] #{result['id']:4d}  {title[:70]}")
            ok += 1
            ingested.append({"id": result["id"], "title": title, "content": content})
        else:
            print(f"  [FALHA]       {title[:70]}")
            fail += 1
    return ok, fail, ingested


def process_file(
    path: Path,
    context_slug: str,
    node_type: str,
    title: Optional[str] = None,
) -> tuple[int, int]:
    ext = path.suffix.upper().lstrip(".")
    print(f"\n  [ARQUIVO] {path.name}  ({ext})")
    chunks = extract_chunks(path)
    if not chunks:
        return 0, 0
    total_chars = sum(len(c["content"]) for c in chunks)
    print(f"  {len(chunks)} chunk(s) | ~{total_chars:,} chars")
    resolved = resolve_type(node_type, path, chunks)
    ok, fail, ingested = ingest_chunks(chunks, context_slug, resolved, title, str(path))

    # Auto-link: mapeia refer�ncias cruzadas entre este n� e os j� existentes no grafo
    if ingested:
        print(f"  [LINK] Mapeando interdepend�ncias...")
        total_links = sum(
            auto_link_node(n["id"], n["content"], n["title"]) for n in ingested
        )
        if total_links:
            print(f"  [LINK] {total_links} aresta(s) de refer�ncia mapeada(s)")
        else:
            print(f"  [LINK] Nenhuma refer�ncia cruzada detectada")

    return ok, fail


def process_dir(
    dir_path: Path,
    context_slug: str,
    node_type: str,
    recursive: bool = False,
) -> tuple[int, int]:
    pattern = "**/*" if recursive else "*"
    files = sorted(
        f for f in dir_path.glob(pattern)
        if f.is_file() and f.suffix.lower() in SUPPORTED_EXTS
    )
    print(f"\n[DIRET�RIO] {dir_path}")
    print(f"  {len(files)} arquivo(s) suportado(s) encontrado(s)")
    total_ok = total_fail = 0
    for f in files:
        # each file classified independently in auto mode
        ok, fail = process_file(f, context_slug, node_type)
        total_ok  += ok
        total_fail += fail
    return total_ok, total_fail


# ---------------------------------------------
# CLI
# ---------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        prog="kare_contexto",
        description="Insere arquivos externos no KARE Context Engine (RAG)",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Exemplos:
  # Ingerir um PDF de arquitetura
  python kare_contexto.py --file visao_arq.pdf --context KARE-programa

  # Ingerir planilha de requisitos com tipo symbol
  python kare_contexto.py --file requisitos.xlsx --context ini-001-checkout-mobile --type symbol

  # Ingerir documento Word com t�tulo customizado
  python kare_contexto.py --file ata_kick-off.docx --context ini-002-onboarding-parceiro --title "Ata Reuni�o Kick-off"

  # Ingerir pasta inteira de documentos recursivamente
  python kare_contexto.py --dir ./docs/arquitetura --context KARE-programa --recursive

  # Ingerir sem reconstruir �ndice (batch de v�rios arquivos)
  python kare_contexto.py --file doc1.pdf --context meu-projeto --no-rebuild
  python kare_contexto.py --file doc2.xlsx --context meu-projeto --no-rebuild
  python kare_contexto.py --file doc3.docx --context meu-projeto  # rebuilda no final

Tipos de n� (--type):
  artifact  (padr�o) � documentos, especifica��es, atas, relat�rios, slides
  symbol              � termos de dom�nio, gloss�rio, entidades de neg�cio
  decision            � decis�es t�cnicas, ADRs, escolhas de arquitetura
  concept             � conceitos gerais explicativos
  context             � projetos ou iniciativas como unidade

Formatos suportados:
  .md .txt .rst  � texto/markdown (nativo)
  .csv .json     � dados estruturados (nativo)
  .pdf           � requer: pip install pdfplumber
  .docx          � requer: pip install python-docx
  .pptx          � requer: pip install python-pptx
  .xlsx .xls     � requer: pip install openpyxl
        """,
    )

    src = parser.add_mutually_exclusive_group(required=True)
    src.add_argument("--file", "-f", metavar="PATH", help="Arquivo a ingerir")
    src.add_argument("--dir",  "-d", metavar="PATH", help="Diret�rio a ingerir")

    parser.add_argument(
        "--context", "-c", required=True, metavar="SLUG",
        help="Slug do contexto/projeto (ex: KARE-programa, ini-001-checkout-mobile)",
    )
    parser.add_argument(
        "--type", "-t", default="auto", choices=NODE_TYPES,
        help="Tipo de n� RAG (padr�o: auto � detectado pelo conte�do)",
    )
    parser.add_argument(
        "--title", metavar="TITLE",
        help="T�tulo personalizado (apenas com --file de 1 chunk)",
    )
    parser.add_argument(
        "--recursive", "-r", action="store_true",
        help="Varrer subpastas (apenas com --dir)",
    )
    parser.add_argument(
        "--no-rebuild", action="store_true",
        help="N�o reconstruir �ndice BM25 ao final (�til para batch)",
    )

    args = parser.parse_args()

    # Banner
    tipo_display = args.type if args.type != "auto" else "auto (detectado por arquivo)"
    print("=" * 64)
    print("  KARE Context Engine � /contexto-rag")
    print(f"  contexto : {args.context}")
    print(f"  tipo     : {tipo_display}")
    print("=" * 64)

    total_ok = total_fail = 0

    if args.file:
        path = Path(args.file).resolve()
        if not path.exists():
            print(f"[ERRO] Arquivo n�o encontrado: {path}", file=sys.stderr)
            sys.exit(1)
        total_ok, total_fail = process_file(path, args.context, args.type, args.title)
    else:
        dir_path = Path(args.dir).resolve()
        if not dir_path.is_dir():
            print(f"[ERRO] Diret�rio n�o encontrado: {dir_path}", file=sys.stderr)
            sys.exit(1)
        total_ok, total_fail = process_dir(dir_path, args.context, args.type, args.recursive)

    # Summary
    print("\n" + "=" * 64)
    status = "[OK]" if total_fail == 0 else "[PARCIAL]" if total_ok > 0 else "[FALHA]"
    print(f"  {status}  {total_ok} n�(s) ingerido(s)  |  {total_fail} falha(s)")

    if total_ok > 0 and not args.no_rebuild:
        _rebuild_index()

    print("=" * 64)


if __name__ == "__main__":
    main()
